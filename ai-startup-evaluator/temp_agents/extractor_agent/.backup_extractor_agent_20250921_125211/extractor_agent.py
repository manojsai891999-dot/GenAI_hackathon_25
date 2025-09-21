import os
import json
import logging
from typing import Dict, List, Optional, Any
from PyPDF2 import PdfReader
from pptx import Presentation
import re
from datetime import datetime

from google.adk.agents import Agent

from models.pydantic_models import (
    ExtractorAgentInput, 
    ExtractorAgentOutput, 
    StartupCreate, 
    FounderCreate
)
from services.gcs_service import gcs_service, get_extracted_data_path

logger = logging.getLogger(__name__)

def extract_text_from_pdf(file_path: str) -> dict:
    """Extract text content from PDF file"""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return {"status": "success", "text": text}
    except Exception as e:
        logger.error(f"Failed to extract text from PDF: {str(e)}")
        return {"status": "error", "error_message": str(e)}

def extract_text_from_pptx(file_path: str) -> dict:
    """Extract text content from PowerPoint file"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return {"status": "success", "text": text}
    except Exception as e:
        logger.error(f"Failed to extract text from PPTX: {str(e)}")
        return {"status": "error", "error_message": str(e)}

def extract_startup_data(text_content: str) -> dict:
    """Extract structured startup information from text content"""
    if not text_content or not text_content.strip():
        return {
            "status": "error",
            "error_message": "No text content provided"
        }
    
    # This would be processed by the ADK agent's LLM
    # For now, return a structured response format
    return {
        "status": "success",
        "extraction_result": {
            "startup": {},
            "founders": [],
            "confidence": 0.5,
            "extracted_fields": [],
            "missing_fields": []
        }
    }

def save_extraction_to_gcs(startup_id: int, extraction_data: dict) -> dict:
    """Save extraction data as backup to GCS"""
    try:
        gcs_path = get_extracted_data_path(startup_id)
        backup_data = {
            "extraction_timestamp": datetime.utcnow().isoformat(),
            "agent_version": "1.0",
            "data": extraction_data
        }
        
        gcs_service.upload_json(backup_data, gcs_path.replace("gs://", "").split("/", 1)[1])
        logger.info(f"Extraction backup saved to: {gcs_path}")
        return {"status": "success", "gcs_path": gcs_path}
    except Exception as e:
        logger.error(f"Failed to save extraction backup: {str(e)}")
        return {"status": "error", "error_message": str(e)}

def validate_startup_data(extracted_data: dict) -> dict:
    """Clean and validate extracted startup data"""
    try:
        startup_data = extracted_data.get("startup", {})
        founders_data = extracted_data.get("founders", [])
        
        # Clean startup data
        if startup_data.get("founded_year"):
            try:
                year = int(startup_data["founded_year"])
                if year < 1900 or year > datetime.now().year:
                    startup_data["founded_year"] = None
            except (ValueError, TypeError):
                startup_data["founded_year"] = None
        
        # Clean financial metrics
        financial_fields = ["revenue", "arr", "burn_rate", "funding_raised", "valuation", "cac", "ltv"]
        for field in financial_fields:
            if startup_data.get(field):
                try:
                    value = float(startup_data[field])
                    startup_data[field] = value if value >= 0 else None
                except (ValueError, TypeError):
                    startup_data[field] = None
        
        # Clean rates (should be between 0 and 1)
        rate_fields = ["churn_rate", "growth_rate"]
        for field in rate_fields:
            if startup_data.get(field):
                try:
                    value = float(startup_data[field])
                    if value > 1:  # Convert percentage to decimal
                        value = value / 100
                    startup_data[field] = value if 0 <= value <= 1 else None
                except (ValueError, TypeError):
                    startup_data[field] = None
        
        # Clean founder data
        for founder in founders_data:
            if founder.get("years_experience"):
                try:
                    years = int(founder["years_experience"])
                    founder["years_experience"] = years if 0 <= years <= 50 else None
                except (ValueError, TypeError):
                    founder["years_experience"] = None
        
        return {
            "status": "success",
            "validated_data": {
                "startup": startup_data,
                "founders": founders_data,
                "confidence": extracted_data.get("confidence", 0.5),
                "extracted_fields": extracted_data.get("extracted_fields", []),
                "missing_fields": extracted_data.get("missing_fields", [])
            }
        }
    except Exception as e:
        return {"status": "error", "error_message": str(e)}

# Create the ExtractorAgent using Google ADK
extractor_agent = Agent(
    name="startup_extractor_agent",
    model="gemini-2.0-flash",
    description="Agent specialized in extracting structured startup information from pitch decks, presentations, and form data",
    instruction="""
    You are an expert startup analyst agent. Your role is to extract structured information from startup pitch decks, presentations, and form data.
    
    When processing documents, extract the following information and return it as a JSON object:
    
    STARTUP INFORMATION:
    - name: Company name
    - sector: Industry/sector (e.g., "FinTech", "HealthTech", "SaaS", "E-commerce")
    - stage: Funding stage (e.g., "Pre-seed", "Seed", "Series A", "Series B")
    - description: Brief company description
    - website: Company website URL
    - location: Company location/headquarters
    - founded_year: Year founded
    - team_size: Number of employees
    
    FINANCIAL METRICS:
    - revenue: Annual revenue in USD
    - arr: Annual Recurring Revenue in USD
    - burn_rate: Monthly burn rate in USD
    - runway_months: Runway in months
    - funding_raised: Total funding raised to date in USD
    - valuation: Current valuation in USD
    
    KEY METRICS:
    - cac: Customer Acquisition Cost in USD
    - ltv: Customer Lifetime Value in USD
    - churn_rate: Monthly churn rate as decimal (e.g., 0.05 for 5%)
    - growth_rate: Monthly growth rate as decimal (e.g., 0.15 for 15%)
    
    FOUNDERS (array of objects):
    - name: Founder name
    - role: Role/title (e.g., "CEO", "CTO", "Co-founder")
    - email: Email address
    - linkedin: LinkedIn profile URL
    - background: Brief background description
    - education: Educational background
    - previous_experience: Previous work experience
    - years_experience: Total years of relevant experience
    
    Always return a JSON object with the structure:
    {
        "startup": { startup fields },
        "founders": [ founder objects ],
        "confidence": 0.85,
        "extracted_fields": ["name", "sector", ...],
        "missing_fields": ["valuation", "cac", ...]
    }
    
    If a field cannot be found, set it to null. Be conservative with confidence scores.
    Use the available tools to extract text from files and validate the data.
    """,
    tools=[
        extract_text_from_pdf,
        extract_text_from_pptx,
        extract_startup_data,
        validate_startup_data,
        save_extraction_to_gcs
    ]
)
